"""
Генератор презентаций на Claude.

Pipeline:
  1. Юзер заполняет форму: тема + аудитория + кол-во слайдов + extra_info
     + опционально загружает картинки
  2. _claude_prompt → AI возвращает СТРОГИЙ JSON со слайдами
  3. _render_html_preview → рендерим HTML-карусель для превью на сайте
  4. _build_pptx → собираем .pptx файл через python-pptx
  5. Скачивание PPTX из /uploads/presentations/

Ценообразование: real_tokens × 7 (margin_pct = 700%) с минимумом
из presentation.min_cost (по умолчанию 50 ₽). Юзер видит примерную
стоимость ДО генерации (мелким шрифтом в UI).
"""
import os
import re
import json
import logging
from datetime import datetime
from io import BytesIO

from server.ai import generate_response
from server.models import PresentationProject

log = logging.getLogger(__name__)

# Лимиты
MIN_SLIDES = 3
MAX_SLIDES = 40
MIN_TOPIC_LEN = 5
MAX_EXTRA_INFO = 15_000

# Цветовые схемы для рендера. Backend применяет к HTML preview + PPTX.
_COLOR_SCHEMES = {
    "dark": {
        "bg": "#1C1C1C", "panel": "#272018", "text": "#f0e6d8",
        "accent": "#ff8c42", "accent2": "#ffb347", "muted": "#a89880",
        "title_color": "#ffb347",
    },
    "light": {
        "bg": "#fafaf7", "panel": "#ffffff", "text": "#1a1a1a",
        "accent": "#ff8c42", "accent2": "#d97706", "muted": "#666",
        "title_color": "#1e293b",
    },
    "corp": {
        "bg": "#0f172a", "panel": "#1e293b", "text": "#e2e8f0",
        "accent": "#3b82f6", "accent2": "#1d4ed8", "muted": "#94a3b8",
        "title_color": "#60a5fa",
    },
    "brand": {
        # Подставляется из бренда: см. _resolve_colors
        "bg": "#1C1C1C", "panel": "#272018", "text": "#f0e6d8",
        "accent": "#ff8c42", "accent2": "#ffb347", "muted": "#a89880",
        "title_color": "#ffb347",
    },
}


def _resolve_colors(scheme: str | None) -> dict:
    return _COLOR_SCHEMES.get(scheme or "dark", _COLOR_SCHEMES["dark"])


def _is_hex(s: str | None) -> bool:
    if not s:
        return False
    s = s.strip()
    if not s.startswith("#"):
        return False
    h = s[1:]
    return len(h) in (3, 6) and all(c in "0123456789abcdefABCDEF" for c in h)


def _resolve_colors_for_project(p: PresentationProject) -> dict:
    """Главная функция выбора цветов: кастомные > бренд клиента > пресет."""
    # 1) Если есть свои hex'ы — приоритет
    if (p.bg_color and _is_hex(p.bg_color)) or (p.accent_color and _is_hex(p.accent_color)):
        return _build_custom_palette(
            bg=p.bg_color or "#1C1C1C",
            text=p.text_color or "#f0e6d8",
            accent=p.accent_color or "#ff8c42",
            title=p.title_color or (p.accent_color or "#ffb347"),
        )
    # 2) Готовый пресет
    return _resolve_colors(p.color_scheme)


def _build_custom_palette(bg: str, text: str, accent: str, title: str) -> dict:
    """Собрать палитру из 4 пользовательских цветов."""
    # accent2 — на 15% светлее accent (для градиентов)
    accent2 = _lighten_hex(accent, 0.15) or accent
    # panel — на 5-7% светлее/темнее bg (контраст для карточек)
    panel = _shift_hex(bg, 0.07)
    muted = _shift_hex(text, -0.35)
    return {
        "bg": bg, "panel": panel or bg, "text": text,
        "accent": accent, "accent2": accent2,
        "muted": muted or text, "title_color": title,
    }


def _hex_to_rgb_tuple(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb_to_hex(r: int, g: int, b: int) -> str:
    return f"#{max(0, min(255, r)):02x}{max(0, min(255, g)):02x}{max(0, min(255, b)):02x}"


def _lighten_hex(hex_color: str, factor: float) -> str:
    """Осветлить (factor>0) или затемнить (factor<0) HEX-цвет."""
    if not _is_hex(hex_color):
        return hex_color
    r, g, b = _hex_to_rgb_tuple(hex_color)
    if factor > 0:
        r = int(r + (255 - r) * factor)
        g = int(g + (255 - g) * factor)
        b = int(b + (255 - b) * factor)
    else:
        r = int(r * (1 + factor))
        g = int(g * (1 + factor))
        b = int(b * (1 + factor))
    return _rgb_to_hex(r, g, b)


def _shift_hex(hex_color: str, factor: float) -> str:
    """Сдвинуть hex в светлую/тёмную сторону относительно его яркости.
    Для тёмных цветов осветляем, для светлых — затемняем."""
    if not _is_hex(hex_color):
        return hex_color
    r, g, b = _hex_to_rgb_tuple(hex_color)
    luminosity = (r + g + b) / 3
    is_dark = luminosity < 128
    f = abs(factor) if is_dark else -abs(factor)
    return _lighten_hex(hex_color, f)


# ── Pre-validation ────────────────────────────────────────────────────────


def validate_presentation(p: PresentationProject) -> None:
    """Проверка ДО списания денег. Кидает ValueError с понятным сообщением."""
    topic = (p.topic or p.name or "").strip()
    if len(topic) < MIN_TOPIC_LEN:
        raise ValueError(f"Тема презентации слишком короткая (минимум {MIN_TOPIC_LEN} символов).")
    sc = int(p.slide_count or 10)
    if sc < MIN_SLIDES or sc > MAX_SLIDES:
        raise ValueError(f"Количество слайдов должно быть от {MIN_SLIDES} до {MAX_SLIDES}.")
    if p.extra_info and len(p.extra_info) > MAX_EXTRA_INFO:
        raise ValueError(f"Доп. информация слишком большая (макс {MAX_EXTRA_INFO} символов).")


# ── Цена ──────────────────────────────────────────────────────────────────


def _claude_rates_kop_per_1k(db=None) -> tuple[float, float]:
    """Ставки Claude в копейках за 1000 токенов (input, output).
    Берём из ModelPricing если есть, иначе fallback по официальным ценам
    Anthropic Sonnet × текущий курс USD/RUB.
    """
    if db is not None:
        try:
            from server.models import ModelPricing
            row = db.query(ModelPricing).filter_by(model_id="claude").first()
            if row and (row.ch_per_1k_input or row.ch_per_1k_output):
                return (float(row.ch_per_1k_input or 0),
                        float(row.ch_per_1k_output or 0))
        except Exception:
            pass
    # Fallback: Sonnet 4 = $3 / $15 за 1M токенов = $0.003 / $0.015 за 1k.
    # При курсе ~95 ₽/$ → 0.285 ₽/1k input, 1.425 ₽/1k output.
    # В копейках: 28.5 / 142.5
    return 28.5, 142.5


def estimate_cost_kop(slide_count: int, extra_info_len: int = 0,
                      images_count: int = 0, has_site: bool = False,
                      db=None) -> tuple[int, int]:
    """Оценка стоимости ДО генерации (для UI «≈ X-Y ₽»).
    Внутри: токены Claude × margin (по умолчанию ×7 — pricing.presentation.margin_pct=700).
    Это маржа сервиса — юзеру в UI не показываем как «×7», только итог.

    Динамика: больше слайдов → больше output → дороже. Картинки дороже из-за vision.
    """
    from server.pricing import get_price as _gp
    in_per_k, out_per_k = _claude_rates_kop_per_1k(db)
    input_t = 1000 + slide_count * 100 + int(extra_info_len * 0.4)
    if images_count:
        input_t += images_count * 1500   # vision-картинки дороже
    if has_site:
        input_t += 800
    output_t = slide_count * 420
    base = (input_t / 1000) * in_per_k + (output_t / 1000) * out_per_k
    margin_pct = int(_gp("presentation.margin_pct", default=700))
    cost = max(int(round(base * margin_pct / 100)), 1)
    return int(cost * 0.85), int(cost * 1.15)


def calc_actual_cost_kop(usage: dict, db) -> int:
    """Реальная стоимость по фактическим токенам Claude × margin
    (presentation.margin_pct=700 по умолчанию).
    Маржа — внутреннее правило сервиса, в UI не светится."""
    from server.pricing import get_price as _gp
    in_per_k, out_per_k = _claude_rates_kop_per_1k(db)
    input_t = int(usage.get("input_tokens", 0) or 0)
    output_t = int(usage.get("output_tokens", 0) or 0)
    base = (input_t / 1000) * in_per_k + (output_t / 1000) * out_per_k
    margin_pct = int(_gp("presentation.margin_pct", default=700))
    return max(int(round(base * margin_pct / 100)), 1)


# ── Claude prompt v1 (JSON-first) ─────────────────────────────────────────


def _claude_prompt(p: PresentationProject, image_urls: list[str],
                    image_descriptions: list[str] | None = None,
                    site_ctx: str = "",
                    custom_charts: list[dict] | None = None) -> str:
    """Промпт. AI возвращает JSON: список слайдов.

    Параметры:
    - image_urls — URL картинок (попадают в image_idx)
    - image_descriptions — что AI vision увидел на картинках (опц.)
    - site_ctx — текст с сайта клиента (для подбора стиля/тона)
    - custom_charts — графики, явно заданные юзером (с данными)
    """
    sc = int(p.slide_count or 10)
    audience = (p.audience or "").strip() or "(не указано)"
    topic = (p.topic or p.name or "").strip()

    parts = [
        "Ты — старший дизайнер презентаций. Сделай структуру слайдов на русском языке.",
        "ВАЖНО: верни СТРОГИЙ JSON. Без markdown-обёртки (```), без комментариев, без HTML.",
        "",
        "=== СТРУКТУРА JSON ===",
        '{',
        '  "title": "Название презентации (1 строка)",',
        '  "subtitle": "Подзаголовок / девиз (опц.)",',
        '  "slides": [',
        '    {',
        '      "type": "title" | "content" | "two_column" | "chart" | "quote" | "section" | "cta",',
        '      "title": "Заголовок слайда",',
        '      "subtitle": "Подзаголовок (опц.)",',
        '      "bullets": ["Пункт 1", "Пункт 2", "..."],          // для type=content',
        '      "left_bullets": [...], "right_bullets": [...],     // для type=two_column',
        '      "left_title": "...", "right_title": "...",         // для type=two_column',
        '      "chart": {                                          // для type=chart',
        '         "kind": "bar" | "line" | "pie",',
        '         "labels": ["1 кв.", "2 кв.", "3 кв.", "4 кв."],   // ПО-РУССКИ, понятно — НЕ Q1/Q2',
        '         "values": [100, 250, 400, 520],',
        '         "x_axis": "Период",                                // подпись оси X (опц., 1-3 слова)',
        '         "y_axis": "Выручка, ₽",                            // подпись оси Y (опц.)',
        '         "caption": "Подпись под графиком"',
        '      },',
        '      "quote": "Цитата",                                  // для type=quote',
        '      "quote_author": "Автор",',
        '      "image_idx": 0,                                     // индекс из image_urls (0-based)',
        '      "speaker_notes": "Заметки для докладчика (опц.)"',
        '    }',
        '  ]',
        '}',
        "",
        f"=== ТРЕБОВАНИЯ ===",
        f"• Сделай РОВНО {sc} слайдов в массиве slides.",
        "• 1-й слайд — обязательно type=title с темой и подзаголовком.",
        f"• Последний слайд — type=cta с конкретным призывом для аудитории «{audience}».",
        "• Если в презентации нужны цифры/динамика — добавь хотя бы один слайд type=chart с реалистичными данными.",
        "• Подписи labels в графиках ВСЕГДА на русском, понятные («1 кв.», «Январь», «2024», «Лидген», «Опт»). Английские «Q1»/«FY24»/«MoM» — ЗАПРЕЩЕНЫ.",
        "• Бизнес-тон, без воды и штампов («команда профессионалов» / «гибкий подход» — запрещены).",
        "• 3-5 буллетов на слайд, каждый — конкретный факт/тезис без пустых слов.",
        "• type=section разделяет смысловые блоки в больших презентациях (>=12 слайдов).",
        "• speaker_notes — короткая подсказка докладчику (1-2 предложения), не дублирует bullets.",
        "",
        f"=== ТЕМА ===",
        topic,
        "",
        f"=== АУДИТОРИЯ ===",
        audience,
    ]
    if p.extra_info:
        parts += [
            "",
            "=== ДОПОЛНИТЕЛЬНАЯ ИНФОРМАЦИЯ (используй ОБЯЗАТЕЛЬНО — это центральный контекст) ===",
            p.extra_info[:MAX_EXTRA_INFO],
        ]
    if image_urls:
        parts += [
            "",
            "=== ДОСТУПНЫЕ ИЗОБРАЖЕНИЯ ===",
            "У тебя есть картинки (по индексам). Распредели их по слайдам через image_idx:",
        ]
        for i, u in enumerate(image_urls):
            desc = ""
            if image_descriptions and i < len(image_descriptions):
                desc = " — " + (image_descriptions[i] or "")[:200]
            parts.append(f"  [{i}] {u[:80]}{desc}")
        parts.append("Используй только релевантные слайды. Для type=title часто кладут лого/обложку.")
        parts.append("Если на картинке есть текст/график/диаграмма — упомяни это в bullets соответствующего слайда.")
    if site_ctx:
        parts += [
            "",
            "=== СТИЛЬ И ТОН КЛИЕНТА (с его сайта — учти при выборе формулировок и тематики) ===",
            site_ctx[:4000],
            "Подбери тон и лексику под ИХ бренд, а не нашу платформу. "
            "Если на сайте корпоративный официальный язык — используй его. "
            "Если стартаповский/энергичный — поддержи стиль.",
        ]
    if custom_charts:
        parts += [
            "",
            "=== ЯВНЫЕ ГРАФИКИ ОТ ПОЛЬЗОВАТЕЛЯ (используй ТОЧНО эти данные) ===",
        ]
        for i, ch in enumerate(custom_charts[:10]):
            kind = ch.get("kind") or "bar"
            labels = ch.get("labels") or []
            values = ch.get("values") or []
            title = ch.get("title") or f"График {i+1}"
            parts.append(f"  [{i}] {title} ({kind}): {list(zip(labels, values))}")
        parts.append("Эти графики — ОБЯЗАТЕЛЬНО включи в презентацию (отдельные slides type=chart с этими данными).")
    parts += [
        "",
        f"Сегодня: {datetime.utcnow().strftime('%d.%m.%Y')}.",
        "ВЕРНИ ТОЛЬКО ВАЛИДНЫЙ JSON. БЕЗ ОБЪЯСНЕНИЙ. БЕЗ ОБЁРТКИ.",
    ]
    return "\n".join(parts)


# ── Vision: краткое описание фото через Claude ─────────────────────────


def describe_image_via_claude(image_url: str, user_api_key: str | None = None) -> str:
    """Прокидываем картинку в Claude vision и просим описать одной строкой:
    что на ней изображено + если есть текст/цифры/график — какие.
    Используется для (1) контекста при генерации слайдов, (2) проверки
    что картинка не мусор. Возвращает '' при ошибке."""
    if not image_url:
        return ""
    full_url = image_url
    if image_url.startswith("/"):
        app_url = os.getenv("APP_URL", "https://aiche.ru").rstrip("/")
        full_url = app_url + image_url
    try:
        # Анторопический формат: content = list[{type:'image', source:{type:'url', url:...}}, ...]
        result = generate_response("claude", [{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "url", "url": full_url}},
                {"type": "text", "text":
                    "Опиши кратко (1 предложение, до 150 символов): что изображено? "
                    "Если есть значимый текст/цифры/график — упомяни. Без пред-обычной фразы «На картинке...»."},
            ],
        }], extra={"max_tokens": 200, "model": "claude-haiku-4"}, user_api_key=user_api_key)
        if isinstance(result, dict):
            return str(result.get("content", "") or "")[:300].strip()
    except Exception as e:
        log.warning(f"[vision] image describe failed: {type(e).__name__}: {e}")
    return ""


def parse_client_site_for_style(url: str) -> tuple[str, dict]:
    """Парсит сайт клиента для:
    1) Текстового контекста (подбор тона, лексики, темы)
    2) Доминирующих цветов (опц., через CSS-сниффинг — пока упрощённо)

    Возвращает (text_ctx, dominant_colors_dict).
    """
    try:
        from server.proposal_builder import parse_client_site as _parse_pp
        ctx = _parse_pp(url)
        return (ctx or ""), {}
    except Exception as e:
        log.warning(f"[client-site] parse failed: {type(e).__name__}: {e}")
        return "", {}


def _parse_json(content: str) -> dict | None:
    if not content:
        return None
    s = content.strip()
    s = re.sub(r"^```(?:json)?\s*", "", s, flags=re.IGNORECASE)
    s = re.sub(r"\s*```\s*$", "", s)
    m = re.search(r"\{[\s\S]*\}", s)
    if not m:
        return None
    try:
        d = json.loads(m.group(0))
        return d if isinstance(d, dict) else None
    except Exception:
        return None


# ── HTML preview ──────────────────────────────────────────────────────────


def _render_html_preview(data: dict, scheme: str) -> str:
    """Совместимость: рендер по имени scheme."""
    return _render_html_preview_inner(data, _resolve_colors(scheme))


def _render_html_preview_inner(data: dict, c: dict) -> str:
    """Карусель слайдов для превью на сайте (с навигацией стрелками + точками)."""
    title = data.get("title", "Презентация")
    subtitle = data.get("subtitle", "")
    slides = data.get("slides") or []

    body_parts = []
    for idx, s in enumerate(slides):
        body_parts.append(_render_slide_html(s, c, idx + 1, len(slides)))

    slides_html = "\n".join(body_parts)

    # Минимальный JS для навигации
    js = """
function pNav(d){
  const slides = document.querySelectorAll('.slide');
  if(!slides.length) return;
  let cur = 0;
  slides.forEach((s,i)=>{ if(s.classList.contains('active')) cur = i; });
  let next = Math.max(0, Math.min(slides.length-1, cur+d));
  slides[cur].classList.remove('active');
  slides[next].classList.add('active');
  document.getElementById('pPos').textContent = (next+1)+'/'+slides.length;
}
function pGo(i){
  document.querySelectorAll('.slide').forEach(s=>s.classList.remove('active'));
  document.querySelectorAll('.slide')[i]?.classList.add('active');
  document.getElementById('pPos').textContent = (i+1)+'/'+document.querySelectorAll('.slide').length;
}
document.addEventListener('keydown', e => {
  if(e.key==='ArrowRight'||e.key===' ') pNav(1);
  else if(e.key==='ArrowLeft') pNav(-1);
});
"""

    html_escape = lambda s: (str(s or "").replace("&","&amp;").replace("<","&lt;")
                              .replace(">","&gt;").replace('"',"&quot;"))

    return f"""<!DOCTYPE html>
<html lang="ru"><head><meta charset="utf-8"/>
<title>{html_escape(title)}</title>
<style>
*{{box-sizing:border-box}}
body{{margin:0;background:{c['bg']};color:{c['text']};font-family:Inter,'Segoe UI',sans-serif;
     min-height:100vh;display:flex;flex-direction:column}}
.deck{{flex:1;display:flex;align-items:center;justify-content:center;padding:24px;
      background:radial-gradient(ellipse at center,{c['panel']} 0%,{c['bg']} 80%)}}
.slide{{display:none;width:100%;max-width:1100px;aspect-ratio:16/9;background:{c['panel']};
       border-radius:18px;padding:48px 64px;box-shadow:0 24px 64px rgba(0,0,0,0.3);
       border:1px solid {c['accent']}25;position:relative;overflow:hidden}}
.slide.active{{display:flex;flex-direction:column}}
.slide h1{{font-size:42px;margin:0 0 16px;color:{c['title_color']};font-weight:800;line-height:1.15}}
.slide h2{{font-size:28px;margin:0 0 14px;color:{c['title_color']};font-weight:700}}
.slide p{{font-size:18px;line-height:1.5;margin:0 0 12px;color:{c['text']}}}
.slide .subtitle{{font-size:20px;color:{c['muted']};margin-bottom:16px;font-weight:500}}
.slide ul{{font-size:18px;line-height:1.6;padding-left:0;list-style:none;margin:0}}
.slide ul li{{position:relative;padding-left:28px;margin-bottom:10px}}
.slide ul li:before{{content:'';position:absolute;left:0;top:10px;width:14px;height:14px;
                     border-radius:4px;background:{c['accent']}}}
.slide.title{{justify-content:center;text-align:center;
              background:linear-gradient(135deg,{c['accent']}15,{c['accent2']}15)}}
.slide.title h1{{font-size:56px;margin-bottom:24px}}
.slide.section{{justify-content:center;text-align:center}}
.slide.section h1{{font-size:64px;color:{c['accent']};text-transform:uppercase;letter-spacing:2px}}
.slide.cta{{justify-content:center;align-items:center;text-align:center;
            background:linear-gradient(135deg,{c['accent']},{c['accent2']});color:#fff}}
.slide.cta h1{{color:#fff;font-size:48px}}
.slide.cta p{{color:#fff;font-size:22px}}
.slide.quote{{justify-content:center;text-align:center}}
.slide.quote .qtext{{font-size:36px;line-height:1.3;font-style:italic;color:{c['accent2']};
                     max-width:850px;margin:0 auto 18px}}
.slide.quote .qauthor{{font-size:18px;color:{c['muted']}}}
.cols{{display:grid;grid-template-columns:1fr 1fr;gap:48px;margin-top:8px}}
.col h3{{font-size:22px;color:{c['accent']};margin:0 0 12px}}
.chart-wrap{{margin-top:14px;flex:1;display:flex;align-items:center;justify-content:center;width:100%}}
.chart-svg{{width:100%;height:340px}}
.slide-img{{max-height:240px;max-width:100%;object-fit:cover;border-radius:10px;margin-top:14px;
            align-self:flex-start}}
.slide.title .slide-img{{max-height:160px;margin:0 auto 24px}}
.slidenum{{position:absolute;bottom:24px;right:32px;font-size:13px;color:{c['muted']};
           letter-spacing:1px;font-weight:600}}
.nav{{padding:14px 24px;display:flex;align-items:center;gap:12px;justify-content:center;
      background:{c['panel']};border-top:1px solid {c['accent']}30}}
.nav button{{background:{c['accent']};color:#fff;border:0;border-radius:10px;padding:10px 20px;
             cursor:pointer;font-weight:700;font-size:14px}}
.nav button:hover{{opacity:0.9}}
.nav .pos{{font-weight:700;font-size:14px;min-width:60px;text-align:center}}
.dots{{display:flex;gap:6px;flex-wrap:wrap;max-width:300px;justify-content:center}}
.dots span{{width:9px;height:9px;border-radius:50%;background:{c['muted']};opacity:0.4;cursor:pointer}}
.dots span.act{{background:{c['accent']};opacity:1}}
.cta-action{{margin-top:14px;font-size:24px;font-weight:700}}
.note{{position:absolute;bottom:60px;left:32px;right:32px;font-size:12px;color:{c['muted']};
       opacity:0.5;font-style:italic}}
</style></head><body>
<div class="deck">
{slides_html}
</div>
<div class="nav">
  <button onclick="pNav(-1)">‹ Назад</button>
  <span class="pos" id="pPos">1/{len(slides)}</span>
  <button onclick="pNav(1)">Вперёд ›</button>
  <div class="dots">
    {''.join(f'<span class="{"act" if i==0 else ""}" onclick="pGo({i})"></span>' for i in range(len(slides)))}
  </div>
</div>
<script>{js}</script>
</body></html>"""


def _render_slide_html(s: dict, c: dict, num: int, total: int) -> str:
    """Один слайд в HTML."""
    typ = (s.get("type") or "content").strip().lower()
    title = (s.get("title") or "").strip()
    subtitle = (s.get("subtitle") or "").strip()
    bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
    image_url = s.get("image_url", "")
    note = (s.get("speaker_notes") or "").strip()
    esc = lambda x: (str(x or "").replace("&","&amp;").replace("<","&lt;")
                      .replace(">","&gt;").replace('"',"&quot;"))
    cls = "slide " + typ + (" active" if num == 1 else "")

    inner = ""
    if typ == "title":
        if image_url:
            inner += f'<img class="slide-img" src="{esc(image_url)}"/>'
        inner += f'<h1>{esc(title)}</h1>'
        if subtitle:
            inner += f'<div class="subtitle">{esc(subtitle)}</div>'
    elif typ == "section":
        inner += f'<h1>{esc(title)}</h1>'
        if subtitle:
            inner += f'<div class="subtitle">{esc(subtitle)}</div>'
    elif typ == "two_column":
        lt = (s.get("left_title") or "").strip()
        rt = (s.get("right_title") or "").strip()
        lb = s.get("left_bullets") if isinstance(s.get("left_bullets"), list) else []
        rb = s.get("right_bullets") if isinstance(s.get("right_bullets"), list) else []
        inner += f'<h2>{esc(title)}</h2>'
        if subtitle:
            inner += f'<div class="subtitle">{esc(subtitle)}</div>'
        inner += '<div class="cols">'
        inner += '<div class="col">'
        if lt: inner += f'<h3>{esc(lt)}</h3>'
        if lb: inner += '<ul>' + ''.join(f'<li>{esc(x)}</li>' for x in lb) + '</ul>'
        inner += '</div><div class="col">'
        if rt: inner += f'<h3>{esc(rt)}</h3>'
        if rb: inner += '<ul>' + ''.join(f'<li>{esc(x)}</li>' for x in rb) + '</ul>'
        inner += '</div></div>'
    elif typ == "chart":
        inner += f'<h2>{esc(title)}</h2>'
        if subtitle:
            inner += f'<div class="subtitle">{esc(subtitle)}</div>'
        chart = s.get("chart") if isinstance(s.get("chart"), dict) else {}
        inner += '<div class="chart-wrap">' + _render_chart_svg(chart, c) + '</div>'
        cap = (chart.get("caption") or "").strip()
        if cap:
            inner += f'<div class="subtitle" style="text-align:center;margin-top:6px">{esc(cap)}</div>'
    elif typ == "quote":
        q = (s.get("quote") or "").strip()
        a = (s.get("quote_author") or "").strip()
        inner += f'<div class="qtext">«{esc(q)}»</div>'
        if a:
            inner += f'<div class="qauthor">— {esc(a)}</div>'
    elif typ == "cta":
        inner += f'<h1>{esc(title)}</h1>'
        if subtitle:
            inner += f'<p style="font-size:24px">{esc(subtitle)}</p>'
        if bullets:
            inner += '<ul style="text-align:left;display:inline-block">' + ''.join(
                f'<li>{esc(x)}</li>' for x in bullets) + '</ul>'
    else:  # content
        inner += f'<h2>{esc(title)}</h2>'
        if subtitle:
            inner += f'<div class="subtitle">{esc(subtitle)}</div>'
        if bullets:
            inner += '<ul>' + ''.join(f'<li>{esc(x)}</li>' for x in bullets) + '</ul>'
        if image_url:
            inner += f'<img class="slide-img" src="{esc(image_url)}"/>'

    if note:
        inner += f'<div class="note" title="Speaker notes">📝 {esc(note)}</div>'

    inner += f'<div class="slidenum">{num} / {total}</div>'
    return f'<div class="{cls}">{inner}</div>'


def _render_chart_svg(chart: dict, c: dict) -> str:
    """Простая SVG-визуализация bar/line/pie. Для PPTX тут же используем."""
    kind = (chart.get("kind") or "bar").lower()
    labels = chart.get("labels") or []
    values = chart.get("values") or []
    if not labels or not values or len(labels) != len(values):
        return '<svg class="chart-svg" viewBox="0 0 600 340"></svg>'
    try:
        values = [float(v) for v in values]
    except Exception:
        return '<svg class="chart-svg" viewBox="0 0 600 340"></svg>'
    max_v = max(values) if values else 1
    if max_v <= 0:
        max_v = 1

    if kind == "bar":
        n = len(labels)
        w = 540 / n
        bars = []
        for i, (lbl, v) in enumerate(zip(labels, values)):
            h = (v / max_v) * 240
            x = 30 + i * w + w * 0.15
            y = 280 - h
            bw = w * 0.7
            bars.append(f'<rect x="{x:.0f}" y="{y:.0f}" width="{bw:.0f}" height="{h:.0f}" '
                        f'fill="{c["accent"]}" rx="4"/>')
            bars.append(f'<text x="{x + bw/2:.0f}" y="{y - 6:.0f}" fill="{c["text"]}" '
                        f'font-size="12" text-anchor="middle" font-weight="600">{v:g}</text>')
            bars.append(f'<text x="{x + bw/2:.0f}" y="305" fill="{c["muted"]}" '
                        f'font-size="11" text-anchor="middle">{_html_safe(lbl)[:14]}</text>')
        return f'<svg class="chart-svg" viewBox="0 0 600 340">{"".join(bars)}</svg>'

    if kind == "line":
        n = len(labels)
        if n < 2:
            return '<svg class="chart-svg" viewBox="0 0 600 340"></svg>'
        step = 540 / (n - 1)
        pts = []
        for i, v in enumerate(values):
            x = 30 + i * step
            y = 280 - (v / max_v) * 240
            pts.append(f"{x:.0f},{y:.0f}")
        polyline = f'<polyline points="{" ".join(pts)}" fill="none" stroke="{c["accent"]}" stroke-width="3"/>'
        dots = "".join(f'<circle cx="{30 + i * step:.0f}" cy="{280 - (v / max_v) * 240:.0f}" '
                       f'r="5" fill="{c["accent2"]}"/>' for i, v in enumerate(values))
        labs = "".join(f'<text x="{30 + i * step:.0f}" y="305" fill="{c["muted"]}" '
                       f'font-size="11" text-anchor="middle">{_html_safe(lbl)[:14]}</text>'
                       for i, lbl in enumerate(labels))
        return f'<svg class="chart-svg" viewBox="0 0 600 340">{polyline}{dots}{labs}</svg>'

    if kind == "pie":
        total = sum(values) or 1
        cx, cy, r = 200, 170, 130
        import math
        a0 = -math.pi / 2
        slices = []
        legend = []
        palette = [c["accent"], c["accent2"], "#ef4444", "#10b981", "#3b82f6", "#a855f7", "#eab308"]
        for i, (lbl, v) in enumerate(zip(labels, values)):
            frac = v / total
            a1 = a0 + frac * 2 * math.pi
            x1 = cx + r * math.cos(a0); y1 = cy + r * math.sin(a0)
            x2 = cx + r * math.cos(a1); y2 = cy + r * math.sin(a1)
            large = 1 if frac > 0.5 else 0
            color = palette[i % len(palette)]
            slices.append(f'<path d="M {cx} {cy} L {x1:.1f} {y1:.1f} '
                          f'A {r} {r} 0 {large} 1 {x2:.1f} {y2:.1f} Z" fill="{color}"/>')
            ly = 60 + i * 26
            legend.append(f'<rect x="380" y="{ly}" width="14" height="14" fill="{color}" rx="2"/>')
            legend.append(f'<text x="402" y="{ly + 11}" fill="{c["text"]}" font-size="13">'
                          f'{_html_safe(lbl)[:22]} — {v:g} ({frac*100:.0f}%)</text>')
            a0 = a1
        return f'<svg class="chart-svg" viewBox="0 0 600 340">{"".join(slices)}{"".join(legend)}</svg>'

    return '<svg class="chart-svg" viewBox="0 0 600 340"></svg>'


def _html_safe(s) -> str:
    return (str(s or "").replace("&","&amp;").replace("<","&lt;").replace(">","&gt;"))


# ── PPTX builder ──────────────────────────────────────────────────────────


def build_pptx(data: dict, scheme: str = "dark", out_path: str = None) -> str:
    """Совместимость: сборка PPTX по имени scheme."""
    return build_pptx_with_palette(data, _resolve_colors(scheme), out_path)


def build_pptx_with_palette(data: dict, c: dict, out_path: str = None) -> str:
    """Сборка .pptx через python-pptx с явной палитрой. Возвращает путь."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    def hex_to_rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    bg = hex_to_rgb(c["bg"])
    panel = hex_to_rgb(c["panel"])
    text_c = hex_to_rgb(c["text"])
    accent = hex_to_rgb(c["accent"])
    title_c = hex_to_rgb(c["title_color"])
    muted = hex_to_rgb(c["muted"])

    title = data.get("title", "Презентация")
    subtitle = data.get("subtitle", "")
    slides = data.get("slides") or []

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        # Фон
        from pptx.dml.color import RGBColor as _RGB  # noqa
        bg_shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)  # 1 = rectangle
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg
        bg_shape.line.fill.background()
        # Без выделения
        bg_shape.shadow.inherit = False

        typ = (s.get("type") or "content").lower()
        st = (s.get("title") or "").strip()
        sub = (s.get("subtitle") or "").strip()

        # Заголовок
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p_ = tf.paragraphs[0]
        p_.text = st
        for run in p_.runs:
            run.font.size = Pt(36 if typ in ("title","section","cta") else 28)
            run.font.bold = True
            run.font.color.rgb = title_c

        if sub:
            sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.0), Inches(0.8))
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = sub
            for run in p2.runs:
                run.font.size = Pt(20 if typ == "title" else 16)
                run.font.color.rgb = muted

        # Тело
        body_top = Inches(2.5)
        body_h = Inches(4.5)

        if typ in ("content", "cta"):
            bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
            if bullets:
                tb = slide.shapes.add_textbox(Inches(0.6), body_top, Inches(12.0), body_h)
                tf = tb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets[:8]):
                    p_ = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p_.text = "•  " + str(b)
                    for run in p_.runs:
                        run.font.size = Pt(18)
                        run.font.color.rgb = text_c

        elif typ == "two_column":
            lb = s.get("left_bullets") or []
            rb = s.get("right_bullets") or []
            lt = (s.get("left_title") or "").strip()
            rt = (s.get("right_title") or "").strip()
            for col_x, col_title, col_bullets in [(0.6, lt, lb), (6.9, rt, rb)]:
                tb = slide.shapes.add_textbox(Inches(col_x), body_top, Inches(5.9), body_h)
                tf = tb.text_frame
                tf.word_wrap = True
                first = True
                if col_title:
                    p_ = tf.paragraphs[0]
                    p_.text = col_title
                    for run in p_.runs:
                        run.font.size = Pt(20); run.font.bold = True
                        run.font.color.rgb = accent
                    first = False
                for b in col_bullets[:6]:
                    p_ = tf.paragraphs[0] if first else tf.add_paragraph()
                    p_.text = "•  " + str(b)
                    for run in p_.runs:
                        run.font.size = Pt(16)
                        run.font.color.rgb = text_c
                    first = False

        elif typ == "quote":
            q = (s.get("quote") or "").strip()
            a = (s.get("quote_author") or "").strip()
            tb = slide.shapes.add_textbox(Inches(1.0), body_top, Inches(11.3), body_h)
            tf = tb.text_frame; tf.word_wrap = True
            p_ = tf.paragraphs[0]
            p_.text = "«" + q + "»"
            for run in p_.runs:
                run.font.size = Pt(28)
                run.font.italic = True
                run.font.color.rgb = title_c
            if a:
                p2 = tf.add_paragraph()
                p2.text = "— " + a
                for run in p2.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = muted

        elif typ == "chart":
            chart = s.get("chart") if isinstance(s.get("chart"), dict) else None
            if chart and chart.get("labels") and chart.get("values"):
                _add_chart_to_slide(slide, chart, body_top, accent)

        # Изображение слайда (если есть и был указан image_url)
        img_url = s.get("image_url") or ""
        if img_url and typ in ("title", "content"):
            try:
                _add_remote_image(slide, img_url, prs)
            except Exception as e:
                log.warning(f"[pptx] image add failed: {type(e).__name__}: {e}")

        # Speaker notes
        notes = (s.get("speaker_notes") or "").strip()
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

        # Номер слайда
        snum = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
        snum.text_frame.text = f"{idx+1} / {len(slides)}"
        for run in snum.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11); run.font.color.rgb = muted

    # Сохранение
    if not out_path:
        out_path = f"/tmp/presentation_{int(datetime.utcnow().timestamp())}.pptx"
    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path


def _add_chart_to_slide(slide, chart_data: dict, top, accent_rgb):
    """Добавляет нативный chart через python-pptx (bar/line)."""
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        kind_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        kind = kind_map.get((chart_data.get("kind") or "bar").lower(),
                             XL_CHART_TYPE.COLUMN_CLUSTERED)
        cd = CategoryChartData()
        cd.categories = [str(x) for x in chart_data.get("labels") or []]
        cd.add_series(chart_data.get("caption") or "Серия 1",
                       [float(v) for v in chart_data.get("values") or []])
        slide.shapes.add_chart(
            kind, Inches(1.5), top, Inches(10.3), Inches(4.5), cd
        )
    except Exception as e:
        log.warning(f"[pptx] chart add failed: {type(e).__name__}: {e}")


def _add_remote_image(slide, url: str, prs) -> None:
    """Скачивает картинку и вставляет в слайд (правый верхний угол)."""
    import httpx
    from pptx.util import Inches
    full_url = url
    if url.startswith("/"):
        app_url = os.getenv("APP_URL", "https://aiche.ru").rstrip("/")
        full_url = app_url + url
    try:
        with httpx.Client(timeout=10) as client:
            r = client.get(full_url)
            if r.status_code != 200 or len(r.content) > 5 * 1024 * 1024:
                return
            buf = BytesIO(r.content)
            slide.shapes.add_picture(buf, Inches(8.5), Inches(2.5),
                                      width=Inches(4.0), height=Inches(4.0))
    except Exception as e:
        log.warning(f"[pptx] download image failed: {type(e).__name__}")


# ── Главный entry point ──────────────────────────────────────────────────


def generate_presentation(db, project: PresentationProject, image_urls: list[str] | None = None,
                           user_api_key: str | None = None) -> dict:
    """Pipeline:
       1. (опц.) Парсим сайт клиента → текст-контекст + цвета
       2. (опц.) Vision-описания всех фото (Claude Haiku)
       3. (опц.) Достаём custom_charts от юзера
       4. AI → JSON со слайдами (с учётом всего контекста)
       5. Маппинг image_idx → image_urls
       6. HTML preview + PPTX

    Возвращает {data, html_path, pptx_path, pdf_path, usage}.
    """
    image_urls = image_urls or []

    # 1) Сайт клиента (если задан)
    site_ctx = ""
    if project.client_site_url and project.client_site_url.strip():
        # Кэш: не парсим повторно если есть client_site_ctx
        if project.client_site_ctx:
            site_ctx = project.client_site_ctx
        else:
            site_ctx, _palette = parse_client_site_for_style(project.client_site_url.strip())
            if site_ctx:
                project.client_site_ctx = site_ctx

    # 2) Vision: описания картинок (только если их немного — экономим)
    image_descriptions: list[str] = []
    if image_urls and len(image_urls) <= 8:
        for url in image_urls:
            desc = describe_image_via_claude(url, user_api_key)
            image_descriptions.append(desc)
            log.info(f"[vision] {url[:60]} → {desc[:80]}")

    # 3) Custom charts от юзера
    custom_charts: list[dict] = []
    if project.custom_charts:
        try:
            cc = json.loads(project.custom_charts)
            if isinstance(cc, list):
                custom_charts = cc[:10]
        except Exception:
            pass

    # 4) Главный prompt → AI
    prompt = _claude_prompt(project, image_urls,
                             image_descriptions=image_descriptions,
                             site_ctx=site_ctx,
                             custom_charts=custom_charts)
    log.info(f"[presentation] gen project={project.id} slides={project.slide_count} "
             f"images={len(image_urls)} site={bool(site_ctx)} charts={len(custom_charts)}")
    ans = generate_response("claude", [{"role": "user", "content": prompt}],
                             extra={"max_tokens": 12000}, user_api_key=user_api_key)
    if not isinstance(ans, dict) or not ans.get("content"):
        raise ValueError("AI вернул пустой ответ")
    data = _parse_json(ans.get("content", ""))
    if not data or not data.get("slides"):
        raise ValueError("AI вернул не-JSON / без слайдов. Попробуйте упростить тему или уменьшить число слайдов.")

    # Маппинг image_idx → image_urls
    for s in data.get("slides", []):
        if isinstance(s, dict) and isinstance(s.get("image_idx"), int):
            idx = s["image_idx"]
            if 0 <= idx < len(image_urls):
                s["image_url"] = image_urls[idx]

    # Если юзер задал custom_charts — ОБЯЗАТЕЛЬНО прокинем их в финальный JSON
    # (страховка на случай если AI забыл их включить)
    if custom_charts:
        existing_charts = sum(1 for s in data.get("slides", [])
                              if isinstance(s, dict) and (s.get("type") or "") == "chart")
        if existing_charts < len(custom_charts):
            # Добавляем недостающие графики после первой content-секции
            insert_pos = 2
            for ch in custom_charts:
                data["slides"].insert(insert_pos, {
                    "type": "chart",
                    "title": ch.get("title") or "Данные",
                    "subtitle": ch.get("subtitle") or "",
                    "chart": {
                        "kind": (ch.get("kind") or "bar").lower(),
                        "labels": ch.get("labels") or [],
                        "values": ch.get("values") or [],
                        "caption": ch.get("caption") or "",
                    },
                })
                insert_pos += 1

    # Палитра — кастомная или пресет
    palette = _resolve_colors_for_project(project)

    # 5) Сохраняем артефакты
    base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    out_dir = os.path.join(base, "uploads", "presentations")
    os.makedirs(out_dir, exist_ok=True)
    ts = int(datetime.utcnow().timestamp())

    # HTML preview
    html = _render_html_preview_with_palette(data, palette)
    html_name = f"pres_{project.id}_{ts}.html"
    with open(os.path.join(out_dir, html_name), "w", encoding="utf-8") as f:
        f.write(html)
    html_rel = f"/uploads/presentations/{html_name}"

    # PPTX
    pptx_rel = None
    pptx_name = f"pres_{project.id}_{ts}.pptx"
    pptx_path_abs = os.path.join(out_dir, pptx_name)
    try:
        build_pptx_with_palette(data, palette, pptx_path_abs)
        pptx_rel = f"/uploads/presentations/{pptx_name}"
    except Exception as e:
        log.error(f"[presentation] PPTX build failed: {type(e).__name__}: {e}")

    # PDF (опц., через xhtml2pdf — может не сработать на сложных слайдах)
    pdf_rel = None
    try:
        from server.pdf_builder import html_to_pdf_bytes
        # Для PDF используем упрощённый «один-слайд-на-страницу» рендер
        pdf_html = _render_pdf_html(data, palette)
        pdf_bytes = html_to_pdf_bytes(pdf_html)
        pdf_name = f"pres_{project.id}_{ts}.pdf"
        with open(os.path.join(out_dir, pdf_name), "wb") as f:
            f.write(pdf_bytes)
        pdf_rel = f"/uploads/presentations/{pdf_name}"
    except Exception as e:
        log.warning(f"[presentation] PDF build skipped: {type(e).__name__}: {e}")

    return {
        "data": data,
        "html_path": html_rel,
        "pptx_path": pptx_rel,
        "pdf_path": pdf_rel,
        "usage": ans.get("usage", {}) or {},
    }


def _render_html_preview_with_palette(data: dict, palette: dict) -> str:
    """Враппер для нового _render_html_preview с явной палитрой
    (вместо имени schemeы)."""
    return _render_html_preview_inner(data, palette)


def _render_pdf_html(data: dict, palette: dict) -> str:
    """HTML-вариант для печати в PDF: каждый слайд на A4-странице."""
    p = palette
    slides = data.get("slides") or []
    parts = [f"""<!DOCTYPE html><html lang="ru"><head><meta charset="utf-8"/>
<style>
@page {{ size: A4 landscape; margin: 0; }}
body {{ margin: 0; font-family: 'Liberation Sans', Inter, sans-serif;
        background: {p['bg']}; color: {p['text']}; }}
.slide {{ page-break-after: always; padding: 32pt 40pt; min-height: 90vh;
          background: {p['bg']}; box-sizing: border-box; }}
h1 {{ color: {p['title_color']}; font-size: 36pt; margin: 0 0 16pt; line-height: 1.2; }}
h2 {{ color: {p['title_color']}; font-size: 26pt; margin: 0 0 12pt; }}
h3 {{ color: {p['accent']}; font-size: 18pt; }}
ul {{ font-size: 16pt; line-height: 1.6; padding-left: 28pt; }}
li {{ margin-bottom: 8pt; }}
.subtitle {{ font-size: 18pt; color: {p['muted']}; }}
.cta {{ background: {p['accent']}; color: #fff; padding: 24pt; border-radius: 6pt; text-align: center; }}
</style></head><body>"""]
    esc = _html_safe
    for s in slides:
        typ = (s.get("type") or "content").lower()
        cls = "slide " + typ
        title = esc(s.get("title", ""))
        sub = esc(s.get("subtitle", ""))
        bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
        parts.append(f'<div class="{cls}">')
        if typ in ("title", "section", "cta"):
            parts.append(f'<h1>{title}</h1>')
        else:
            parts.append(f'<h2>{title}</h2>')
        if sub:
            parts.append(f'<div class="subtitle">{sub}</div>')
        if bullets:
            parts.append('<ul>' + "".join(f'<li>{esc(b)}</li>' for b in bullets) + '</ul>')
        parts.append('</div>')
    parts.append('</body></html>')
    return "".join(parts)
